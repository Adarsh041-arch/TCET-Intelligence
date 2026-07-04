"""Extract features (fonts, colors, styles, structure) from uploaded documents.

Used by the documentation agent to replicate a reference document's style.
"""
import os
from typing import Dict, Any, Optional
from pathlib import Path


def extract_features(filepath: str) -> Dict[str, Any]:
    ext = os.path.splitext(filepath)[1].lower()
    extractors = {
        ".docx": _extract_docx,
        ".pdf": _extract_pdf,
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
    }
    extractor = extractors.get(ext, _extract_unknown)
    try:
        return extractor(filepath)
    except Exception as e:
        return {"format": ext, "error": str(e), "features": {}}


def _extract_docx(filepath: str) -> Dict[str, Any]:
    from docx import Document
    doc = Document(filepath)
    features = {
        "format": "DOCX",
        "paragraphs": len(doc.paragraphs),
        "tables": len(doc.tables),
        "sections": len(doc.sections),
        "styles_used": [],
        "fonts_detected": set(),
        "colors_detected": set(),
        "has_headings": False,
        "heading_levels": set(),
        "has_lists": False,
        "has_tables": False,
        "has_images": False,
        "structure": [],
    }
    for p in doc.paragraphs:
        if p.style:
            features["styles_used"].append(p.style.name)
        if p.style and p.style.name.startswith("Heading"):
            features["has_headings"] = True
            try:
                features["heading_levels"].add(int(p.style.name.split()[-1]))
            except (ValueError, IndexError):
                features["heading_levels"].add(1)
        for run in p.runs:
            if run.font.name:
                features["fonts_detected"].add(run.font.name)
            if run.font.color and run.font.color.rgb:
                features["colors_detected"].add(str(run.font.color.rgb))
    if doc.tables:
        features["has_tables"] = True
        features["table_count"] = len(doc.tables)
    features["fonts_detected"] = list(features["fonts_detected"])[:5]
    features["colors_detected"] = list(features["colors_detected"])[:5]
    features["styles_used"] = list(set(features["styles_used"]))[:10]
    features["heading_levels"] = sorted(features["heading_levels"])
    return features


def _extract_pdf(filepath: str) -> Dict[str, Any]:
    features = {
        "format": "PDF",
        "pages": 0,
        "has_text": False,
        "has_tables": False,
        "has_images": False,
    }
    try:
        import PyPDF2
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            features["pages"] = len(reader.pages)
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    features["has_text"] = True
                    break
    except ImportError:
        try:
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                features["pages"] = len(pdf.pages)
                features["has_text"] = any(page.extract_text() for page in pdf.pages if page.extract_text())
        except ImportError:
            features["error"] = "No PDF parser available (install PyPDF2 or pdfplumber)"
    return features


def _extract_pptx(filepath: str) -> Dict[str, Any]:
    from pptx import Presentation
    prs = Presentation(filepath)
    features = {
        "format": "PPTX",
        "slides": len(prs.slides),
        "slide_width": str(prs.slide_width),
        "slide_height": str(prs.slide_height),
        "has_text": False,
        "has_tables": False,
        "has_images": False,
        "fonts_detected": set(),
    }
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                features["has_text"] = True
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.name:
                            features["fonts_detected"].add(run.font.name)
            if shape.has_table:
                features["has_tables"] = True
            if hasattr(shape, "image"):
                features["has_images"] = True
    features["fonts_detected"] = list(features["fonts_detected"])[:5]
    return features


def _extract_xlsx(filepath: str) -> Dict[str, Any]:
    from openpyxl import load_workbook
    wb = load_workbook(filepath, data_only=True)
    features = {
        "format": "XLSX",
        "sheets": len(wb.sheetnames),
        "sheet_names": wb.sheetnames[:5],
        "has_formulas": False,
        "has_tables": False,
        "total_rows": 0,
    }
    for name in wb.sheetnames[:3]:
        ws = wb[name]
        features["total_rows"] += ws.max_row or 0
        if ws.tables:
            features["has_tables"] = True
    return features


def _extract_unknown(filepath: str) -> Dict[str, Any]:
    return {"format": os.path.splitext(filepath)[1].lower(), "error": "Unsupported format for feature extraction"}
