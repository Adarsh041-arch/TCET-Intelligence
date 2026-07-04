import io
import re
import math
from typing import Dict, Any, Optional, List
from html.parser import HTMLParser
from app.document_generation.generators.base import BaseGenerator


class _PPTXContentParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.slides = []
        self.current_slide = {"title": "", "content": [], "type": "content"}
        self.current_item = ""
        self.in_list = False
        self.list_items = []
        self.in_table = False
        self.table_data = []
        self.current_row = []
        self.in_header = False
        self.in_bold = False
        self.in_italic = False
        self.in_code = False
        self.code_content = ""
        self.in_blockquote = False
        self.heading_count = 0
        self.collect_text = False
        self.text_buffer = ""

    def handle_starttag(self, tag, attrs):
        tag = tag.lower()
        if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            self._flush_text()
            level = int(tag[1])
            if level <= 2:
                if self.current_slide["content"] or self.list_items:
                    self._finalize_slide()
                self.current_slide = {"title": "", "content": [], "type": "title"}
            self.collect_text = True
            self.text_buffer = ""
        elif tag == "p":
            self.collect_text = True
            self.text_buffer = ""
        elif tag in ("ul", "ol"):
            self.in_list = True
            self.list_items = []
        elif tag == "li":
            if self.text_buffer:
                self._flush_text()
        elif tag in ("strong", "b"):
            self.in_bold = True
        elif tag in ("em", "i"):
            self.in_italic = True
        elif tag == "pre":
            self.in_code = True
            self.code_content = ""
        elif tag == "code":
            self.in_code = True
            self.code_content = ""
        elif tag == "blockquote":
            self.in_blockquote = True
        elif tag == "table":
            self.in_table = True
            self.table_data = []
        elif tag == "tr":
            self.current_row = []
        elif tag in ("th", "td"):
            pass

    def handle_endtag(self, tag):
        tag = tag.lower()
        if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            self._flush_text()
            self.collect_text = False
        elif tag == "p":
            self._flush_text()
            self.collect_text = False
        elif tag in ("ul", "ol"):
            if self.list_items:
                self.current_slide["content"].append({"type": "list", "items": list(self.list_items)})
                self.list_items = []
            self.in_list = False
        elif tag == "li":
            self._flush_text()
        elif tag in ("strong", "b"):
            self.in_bold = False
        elif tag in ("em", "i"):
            self.in_italic = False
        elif tag == "pre":
            if self.code_content:
                self.current_slide["content"].append({"type": "code", "text": self.code_content})
            self.in_code = False
        elif tag == "code":
            if self.code_content:
                self.current_slide["content"].append({"type": "code", "text": self.code_content})
            self.in_code = False
        elif tag == "blockquote":
            if self.text_buffer:
                self._flush_text()
            self.in_blockquote = False
        elif tag == "table":
            if self.table_data:
                self.current_slide["content"].append({"type": "table", "data": self.table_data})
            self.in_table = False
        elif tag == "tr":
            if self.current_row:
                self.table_data.append(list(self.current_row))
                self.current_row = []
        elif tag in ("th", "td"):
            pass
        elif tag == "hr":
            self._finalize_slide()
            self.current_slide = {"title": "", "content": [], "type": "separator"}

    def handle_data(self, data):
        if self.in_code:
            self.code_content += data
            return
        if self.in_table:
            self.current_row.append(data.strip())
            return
        if self.in_list:
            if data.strip():
                self.list_items.append(data.strip())
            return
        if self.collect_text:
            self.text_buffer += data
        elif self.current_item:
            self.current_item += data

    def _flush_text(self):
        text = self.text_buffer.strip()
        if text:
            if self.in_list:
                self.list_items.append(text)
            else:
                item = {"type": "text", "value": text}
                if self.in_bold:
                    item["bold"] = True
                if self.in_italic:
                    item["italic"] = True
                if self.in_blockquote:
                    item["quote"] = True
                self.current_slide["content"].append(item)
        self.text_buffer = ""

    def _finalize_slide(self):
        if self.current_slide["title"] or self.current_slide["content"]:
            self.slides.append(dict(self.current_slide))
        self.current_slide = {"title": "", "content": [], "type": "content"}

    def get_slides(self):
        self._flush_text()
        self._finalize_slide()
        return self.slides if self.slides else [{"title": "Content", "content": [{"type": "text", "value": "No content"}], "type": "content"}]


class PPTXGenerator(BaseGenerator):
    @property
    def format(self) -> str:
        return "pptx"

    @property
    def mime_type(self) -> str:
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    @property
    def file_extension(self) -> str:
        return ".pptx"

    def generate(self, html: str, template: Optional[Dict[str, Any]] = None, metadata: Optional[Dict[str, Any]] = None) -> bytes:
        try:
            return self._generate_with_python_pptx(html, template)
        except ImportError:
            raise ImportError("python-pptx is required for PPTX generation. Install with: pip install python-pptx")

    def generate_preview(self, html: str, template: Optional[Dict[str, Any]] = None) -> List[bytes]:
        return self._generate_slide_thumbnails(html)

    def _generate_with_python_pptx(self, html: str, template: Optional[Dict] = None) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor as PptRGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        parser = _PPTXContentParser()
        parser.feed(html)
        slides_data = parser.get_slides()

        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = PptRGBColor(0xFF, 0xFF, 0xFF)

            left = Inches(0.75)
            top = Inches(0.5)
            width = prs.slide_width - Inches(1.5)
            body_top = Inches(1.5)

            if slide_data.get("title"):
                title_box = slide.shapes.add_textbox(left, top, width, Inches(1))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_data["title"]
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = PptRGBColor(0x1A, 0x1A, 0x2E)

            y_pos = body_top
            content = slide_data.get("content", [])
            for item in content:
                if y_pos > Inches(6.5):
                    break

                if item["type"] == "text":
                    text_box = slide.shapes.add_textbox(left, y_pos, width, Inches(0.5))
                    tf = text_box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = item.get("value", "")
                    p.font.size = Pt(18)
                    p.font.color.rgb = PptRGBColor(0x33, 0x33, 0x33)
                    if item.get("bold"):
                        p.font.bold = True
                    if item.get("italic"):
                        p.font.italic = True
                    if item.get("quote"):
                        p.font.italic = True
                        p.font.color.rgb = PptRGBColor(0x66, 0x66, 0x66)
                    y_pos += Inches(0.5)

                elif item["type"] == "list":
                    for li in item.get("items", []):
                        if y_pos > Inches(6.5):
                            break
                        text_box = slide.shapes.add_textbox(left + Inches(0.3), y_pos, width - Inches(0.3), Inches(0.4))
                        tf = text_box.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = f"• {li}"
                        p.font.size = Pt(16)
                        p.font.color.rgb = PptRGBColor(0x33, 0x33, 0x33)
                        y_pos += Inches(0.4)

                elif item["type"] == "code":
                    text_box = slide.shapes.add_textbox(left, y_pos, width, Inches(1.5))
                    tf = text_box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = item.get("text", "")
                    p.font.size = Pt(12)
                    p.font.name = "Courier New"
                    p.font.color.rgb = PptRGBColor(0x2D, 0x2D, 0x2D)
                    y_pos += Inches(1.5)

                elif item["type"] == "table":
                    try:
                        table_data = item.get("data", [])
                        if table_data:
                            rows = len(table_data)
                            cols = max(len(r) for r in table_data)
                            table_shape = slide.shapes.add_table(rows, cols, int(left), int(y_pos), int(width), int(Inches(0.4 * rows)))
                            table = table_shape.table
                            for ri, row_data in enumerate(table_data):
                                for ci, cell_text in enumerate(row_data):
                                    if ci < cols:
                                        cell = table.cell(ri, ci)
                                        cell.text = str(cell_text)
                            y_pos += Inches(0.4 * rows + 0.3)
                    except Exception:
                        pass

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _generate_slide_thumbnails(self, html: str) -> List[bytes]:
        try:
            from pptx import Presentation
            from pptx.util import Inches
            import io
            import tempfile
            import os

            pptx_bytes = self._generate_with_python_pptx(html)
            pngs = []

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(pptx_bytes)
                tmp_path = tmp.name

            try:
                from PIL import Image
                prs = Presentation(tmp_path)
                for i, slide in enumerate(prs.slides):
                    img = Image.new("RGB", (400, 225), (255, 255, 255))
                    buf = io.BytesIO()
                    img.save(buf, format="PNG")
                    pngs.append(buf.getvalue())
            finally:
                os.unlink(tmp_path)

            return pngs if pngs else [self._generate_placeholder_thumb()]
        except ImportError:
            return [self._generate_placeholder_thumb()]

    def _generate_placeholder_thumb(self) -> bytes:
        try:
            from PIL import Image, ImageDraw
            img = Image.new("RGB", (400, 225), (245, 245, 245))
            draw = ImageDraw.Draw(img)
            draw.text((150, 100), "Slide Preview", fill=(100, 100, 100))
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()
        except ImportError:
            import struct
            width, height = 400, 225
            raw = b""
            for y in range(height):
                for x in range(width):
                    raw += b"\xf5\xf5\xf5\xff"
            raw = b"\x89PNG\r\n\x1a\n" + struct.pack(">I", 13) + b"IHDR" + struct.pack(">IIBB", width, height, 8, 6) + b"\x00\x00\x00\x00" + b"IDAT" + raw + b"\x00\x00\x00\x00IEND\xaeB`\x82"
            import zlib
            return raw


def pptx_generate(html: str, template: Optional[Dict] = None, metadata: Optional[Dict] = None) -> bytes:
    gen = PPTXGenerator()
    return gen.generate(html, template, metadata)
