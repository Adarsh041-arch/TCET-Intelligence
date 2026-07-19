import io
from typing import Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN

from app.document_generation.generators.base import BaseGenerator
from app.document_generation.markdown_ast import parse, extract_text


def _split_slides(blocks: list) -> list[list[dict]]:
    """Split AST blocks into slides.

    Heuristic: each H1/H2 starts a new slide. A `---` marker also splits.
    """
    slides = []
    current = []
    for block in blocks:
        if block["type"] == "thematic_break":
            if current:
                slides.append(current)
                current = []
            continue
        if block["type"] == "heading" and block.get("level", 1) <= 2:
            if current:
                slides.append(current)
                current = []
        current.append(block)
    if current:
        slides.append(current)
    if not slides:
        slides.append([])
    return slides


class PPTXGeneratorV2(BaseGenerator):
    @property
    def format(self) -> str:
        return "pptx-v2"

    @property
    def mime_type(self) -> str:
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    @property
    def file_extension(self) -> str:
        return ".pptx"

    def generate(self, html: str, template: Optional[Dict[str, Any]] = None, metadata: Optional[Dict[str, Any]] = None) -> bytes:
        return self._generate_v2(html, template, metadata)

    def generate_preview(self, html: str, template: Optional[Dict[str, Any]] = None) -> list[bytes]:
        try:
            import tempfile
            import os
            pptx_bytes = self.generate(html, template)
            pngs = []
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(pptx_bytes)
                tmp_path = tmp.name
            try:
                from PIL import Image
                prs = Presentation(tmp_path)
                for i, slide in enumerate(prs.slides):
                    img = Image.new("RGB", (400, 225), (255, 255, 255))
                    buf = io.BytesIO()
                    img.save(buf, format="PNG")
                    pngs.append(buf.getvalue())
            finally:
                os.unlink(tmp_path)
            return pngs if pngs else [self._generate_placeholder_thumb()]
        except ImportError:
            return [self._generate_placeholder_thumb()]

    def _generate_placeholder_thumb(self) -> bytes:
        try:
            from PIL import Image, ImageDraw
            img = Image.new("RGB", (400, 225), (245, 245, 245))
            draw = ImageDraw.Draw(img)
            draw.text((150, 100), "Slide Preview", fill=(100, 100, 100))
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()
        except ImportError:
            import struct
            width, height = 400, 225
            raw = b""
            for y in range(height):
                for x in range(width):
                    raw += b"\xf5\xf5\xf5\xff"
            raw = b"\x89PNG\r\n\x1a\n" + struct.pack(">I", 13) + b"IHDR" + struct.pack(">IIBB", width, height, 8, 6) + b"\x00\x00\x00\x00" + b"IDAT" + raw + b"\x00\x00\x00\x00IEND\xaeB`\x82"
            return raw

    def _generate_v2(self, html: str, template: Optional[Dict] = None, metadata: Optional[Dict] = None) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        ast = parse(html)
        slides = _split_slides(ast)

        for slide_blocks in slides:
            self._render_slide(prs, slide_blocks, template)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _render_slide(self, prs: Presentation, blocks: list, template: Optional[Dict] = None):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PptRGBColor(0xFF, 0xFF, 0xFF)

        left = Inches(0.75)
        width = prs.slide_width - Inches(1.5)
        y_pos = Inches(0.5)

        for block in blocks:
            if y_pos > Inches(6.5):
                break
            if block["type"] == "heading":
                text = extract_text(block.get("children", []))
                level = block.get("level", 1)
                font_size = {1: 36, 2: 30, 3: 24, 4: 20, 5: 18, 6: 16}.get(level, 24)
                tb = slide.shapes.add_textbox(left, y_pos, width, Inches(0.8))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(font_size)
                p.font.bold = True
                p.font.color.rgb = PptRGBColor(0x1A, 0x1A, 0x2E)
                y_pos += Inches(0.9)
            elif block["type"] == "paragraph":
                children = block.get("children", [])
                if not children:
                    continue
                h = self._estimate_height(children)
                tb = slide.shapes.add_textbox(left, y_pos, width, h)
                tf = tb.text_frame
                tf.word_wrap = True
                self._render_inline_pptx(tf, children)
                y_pos += h + Inches(0.1)
            elif block["type"] == "list":
                h = self._estimate_list_height(block)
                tb = slide.shapes.add_textbox(left, y_pos, width, h)
                tf = tb.text_frame
                tf.word_wrap = True
                self._render_list_pptx(tf, block)
                y_pos += h + Inches(0.1)
            elif block["type"] == "block_code":
                code = block.get("raw", "")
                tb = slide.shapes.add_textbox(left, y_pos, width, Inches(1.5))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = code
                p.font.size = Pt(12)
                p.font.name = "Courier New"
                p.font.color.rgb = PptRGBColor(0x2D, 0x2D, 0x2D)
                y_pos += Inches(1.5)
            elif block["type"] == "table":
                self._render_table_pptx(slide, block, left, y_pos, width)
                y_pos += Inches(2.0)
            elif block["type"] == "block_quote":
                for child in block.get("children", []):
                    if child["type"] == "paragraph":
                        children = child.get("children", [])
                        if not children:
                            continue
                        tb = slide.shapes.add_textbox(left + Inches(0.3), y_pos, width - Inches(0.3), Inches(0.5))
                        tf = tb.text_frame
                        tf.word_wrap = True
                        self._render_inline_pptx(tf, children, is_quote=True)
                        y_pos += Inches(0.5)

    def _render_inline_pptx(self, tf, children: list, is_quote: bool = False):
        text = self._flatten_inline(children)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.color.rgb = PptRGBColor(0x33, 0x33, 0x33)
        if is_quote:
            p.font.italic = True
            p.font.color.rgb = PptRGBColor(0x66, 0x66, 0x66)

    def _flatten_inline(self, children: list) -> str:
        parts = []
        for child in children:
            t = child["type"]
            if t == "text":
                parts.append(child.get("raw", ""))
            elif t == "codespan":
                parts.append(child.get("raw", ""))
            elif t in ("strong", "emphasis"):
                parts.append(self._flatten_inline(child.get("children", [])))
            elif t == "softbreak":
                parts.append(" ")
            elif t == "linebreak":
                parts.append("\n")
        return "".join(parts)

    def _render_list_pptx(self, tf, block: dict, depth: int = 0):
        first = True
        for item in block.get("children", []):
            children = item.get("children", [])
            text_parts = []
            sub_list = None
            for c in children:
                if c["type"] == "list":
                    sub_list = c
                else:
                    text_parts.append(c)
            if not first:
                tf.add_paragraph()
            first = False
            p = tf.paragraphs[-1]
            text = self._flatten_inline(text_parts)
            prefix = "  " * depth + ("• " if not block.get("ordered") else f"{depth + 1}. ")
            p.text = prefix + text
            p.font.size = Pt(16)
            p.font.color.rgb = PptRGBColor(0x33, 0x33, 0x33)
            p.level = depth
            if sub_list:
                self._render_list_pptx(tf, sub_list, depth + 1)

    def _render_table_pptx(self, slide, block: dict, left, y_pos, width):
        rows_data = []
        for child in block.get("children", []):
            for row in child.get("children", []):
                cells = [extract_text(c.get("children", [])) for c in row.get("children", [])]
                rows_data.append(cells)
        if not rows_data:
            return
        rows = len(rows_data)
        cols = max(len(r) for r in rows_data)
        table_shape = slide.shapes.add_table(rows, cols, int(left), int(y_pos), int(width), int(Inches(0.4 * rows)))
        table = table_shape.table
        for ri, row_data in enumerate(rows_data):
            for ci, cell_text in enumerate(row_data):
                if ci < cols:
                    cell = table.cell(ri, ci)
                    cell.text = str(cell_text)

    def _estimate_height(self, children: list) -> Inches:
        return Inches(0.5)

    def _estimate_list_height(self, block: dict) -> Inches:
        count = len(block.get("children", []))
        return Inches(max(0.4, count * 0.35))


def pptx_generate_v2(html: str, template: Optional[Dict] = None, metadata: Optional[Dict] = None) -> bytes:
    gen = PPTXGeneratorV2()
    return gen.generate(html, template, metadata)
